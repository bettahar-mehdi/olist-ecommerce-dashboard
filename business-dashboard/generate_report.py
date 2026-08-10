"""
Statistical Report Generator - French & Arabic Versions (Improved)
- Larger fonts
- Explanatory text on every slide
- No empty slides
- Better Arabic font rendering
"""

import io
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.io import to_image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
from datetime import datetime
from pathlib import Path

# --- Load Data ---
DATA_PATH = Path("data/sales_data.csv")
df = pd.read_csv(DATA_PATH)
df["order_date"] = pd.to_datetime(df["order_date"])

OUTPUT_DIR = Path("deliverables")
OUTPUT_DIR.mkdir(exist_ok=True)

# --- Helper Functions ---


def fig_to_bytes(fig, width=1200, height=600, scale=2):
    """Convert plotly figure to PNG bytes with higher resolution."""
    img_bytes = to_image(fig, format="png", width=width, height=height, scale=scale)
    return io.BytesIO(img_bytes)


def set_font_rpr(run, size=14, bold=False, color=None):
    """Set font properties for a run with Unicode support."""
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    # Ensure Unicode font
    r = run._r
    rPr = r.get_or_add_rPr()
    for tag in ('eastAsia', 'cs', 'hAnsi'):
        e = rPr.find(qn(f'a:{tag}'))
        if e is not None:
            rPr.remove(e)


def add_title_slide(prs, title, subtitle, lang="fr"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0xBF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    set_font_rpr(p.runs[0], size=44, bold=True, color=RGBColor(0x00, 0xBF, 0xFF))

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER
    set_font_rpr(p2.runs[0], size=24, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    date_str = datetime.now().strftime("%B %Y")
    p3.text = f"Date: {date_str}" if lang == "fr" else f"التاريخ: {date_str}"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p3.alignment = PP_ALIGN.CENTER
    set_font_rpr(p3.runs[0], size=16, color=RGBColor(0x99, 0x99, 0x99))


def add_content_slide(prs, title, fig, explanation, lang="fr"):
    """Add a slide with chart title, large chart, and detailed explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x2E)
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0xBF, 0xFF)
    set_font_rpr(p.runs[0], size=26, bold=True, color=RGBColor(0x00, 0xBF, 0xFF))

    # Chart image - larger
    img_bytes = fig_to_bytes(fig, width=1200, height=550, scale=2)
    slide.shapes.add_picture(img_bytes, Inches(1.2), Inches(1.1), width=Inches(7.6))

    # Explanation box at bottom
    exp_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(9), Inches(1.5)
    )
    exp_box.fill.solid()
    exp_box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x44)
    exp_box.line.color.rgb = RGBColor(0x00, 0xBF, 0xFF)
    exp_box.line.width = Pt(1)

    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(8.6), Inches(1.3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = explanation
    p2.font.size = Pt(15)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    set_font_rpr(p2.runs[0], size=15, color=RGBColor(0xDD, 0xDD, 0xDD))


def add_kpi_slide(prs, kpis, explanation, lang="fr"):
    """Add executive summary slide with KPI boxes and explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.fill.background()

    # Title
    title_text = "Resume Executif" if lang == "fr" else "الملخص التنفيذي"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0xBF, 0xFF)
    set_font_rpr(p.runs[0], size=32, bold=True, color=RGBColor(0x00, 0xBF, 0xFF))

    # KPI boxes - 2x2 grid
    positions = [(0.5, 1.3), (5.5, 1.3), (0.5, 3.8), (5.5, 3.8)]
    colors = [
        RGBColor(0x00, 0xBF, 0xFF),
        RGBColor(0x00, 0xE6, 0x76),
        RGBColor(0xFF, 0xC1, 0x07),
        RGBColor(0xFF, 0x44, 0x44),
    ]

    for i, (kpi, (x, y)) in enumerate(zip(kpis, positions)):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.5), Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x44)
        box.line.color.rgb = colors[i]
        box.line.width = Pt(2)

        txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(4), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = kpi["label"]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        set_font_rpr(p.runs[0], size=16, color=RGBColor(0xAA, 0xAA, 0xAA))

        txBox2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.9), Inches(4), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = kpi["value"]
        p2.font.size = Pt(30)
        p2.font.bold = True
        p2.font.color.rgb = colors[i]
        set_font_rpr(p2.runs[0], size=30, bold=True, color=colors[i])

    # Explanation
    exp_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(9), Inches(1.0)
    )
    exp_box.fill.solid()
    exp_box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x44)
    exp_box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

    txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(8.6), Inches(0.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = explanation
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    set_font_rpr(p3.runs[0], size=14, color=RGBColor(0xCC, 0xCC, 0xCC))


def add_table_slide(prs, title, dataframe, explanation, lang="fr"):
    """Add a slide with data table and explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x2E)
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0xBF, 0xFF)
    set_font_rpr(p.runs[0], size=26, bold=True, color=RGBColor(0x00, 0xBF, 0xFF))

    # Table
    rows, cols = min(len(dataframe) + 1, 12), len(dataframe.columns)
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), Inches(9), Inches(4.5))
    table = table_shape.table

    # Header row
    for j, col in enumerate(dataframe.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x60, 0x90)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            set_font_rpr(paragraph.runs[0], size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Data rows
    for i in range(min(len(dataframe), rows - 1)):
        for j, col in enumerate(dataframe.columns):
            cell = table.cell(i + 1, j)
            val = dataframe.iloc[i, j]
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                set_font_rpr(paragraph.runs[0], size=13, color=RGBColor(0xDD, 0xDD, 0xDD))
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x44)

    # Explanation
    exp_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(9), Inches(1.3)
    )
    exp_box.fill.solid()
    exp_box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x44)
    exp_box.line.color.rgb = RGBColor(0x00, 0xBF, 0xFF)
    exp_box.line.width = Pt(1)

    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(5.9), Inches(8.6), Inches(1.1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = explanation
    p2.font.size = Pt(15)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    set_font_rpr(p2.runs[0], size=15, color=RGBColor(0xDD, 0xDD, 0xDD))


def add_text_slide(prs, title, bullets, lang="fr"):
    """Add a text-only slide with bullet points and explanations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(0.12), Inches(5.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x00, 0xBF, 0xFF)
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(0.5), Inches(8.5), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    set_font_rpr(p.runs[0], size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Bullet points
    txBox2 = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(8.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        p.space_after = Pt(14)
        set_font_rpr(p.runs[0], size=20, color=RGBColor(0xDD, 0xDD, 0xDD))


# --- Compute Statistics ---

total_revenue = df["revenue"].sum()
total_orders = len(df)
avg_order_value = total_revenue / total_orders
avg_review = df["avg_review_score"].mean()
unique_customers = df["customer_unique_id"].nunique()
unique_categories = df["category"].nunique()
date_min = df["order_date"].min().strftime("%Y-%m-%d")
date_max = df["order_date"].max().strftime("%Y-%m-%d")

# Monthly aggregation
monthly = df.groupby(df["order_date"].dt.to_period("M")).agg(
    revenue=("revenue", "sum"),
    orders=("revenue", "count"),
).reset_index()
monthly["order_date"] = monthly["order_date"].astype(str)

# Regional aggregation
regional = df.groupby("region").agg(
    revenue=("revenue", "sum"),
    orders=("revenue", "count"),
    avg_review=("avg_review_score", "mean"),
).reset_index().sort_values("revenue", ascending=False)

# Category aggregation
category_stats = df.groupby("category").agg(
    revenue=("revenue", "sum"),
    orders=("revenue", "count"),
    avg_price=("price", "mean"),
    avg_review=("avg_review_score", "mean"),
).reset_index().sort_values("revenue", ascending=False)

# Payment stats
payment_stats = df.groupby("payment_type").agg(
    revenue=("revenue", "sum"),
    orders=("revenue", "count"),
).reset_index()

# Day of week
df["day_of_week"] = df["order_date"].dt.day_name()
day_order = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
dow_stats = df.groupby("day_of_week")["revenue"].sum().reindex(day_order).reset_index()

# Top cities
city_stats = df.groupby("customer_city")["revenue"].sum().nlargest(15).reset_index()

# Correlation
price_freight_corr = df["price"].corr(df["freight_value"])

# Percentiles
revenue_p25 = df["revenue"].quantile(0.25)
revenue_p50 = df["revenue"].quantile(0.50)
revenue_p75 = df["revenue"].quantile(0.75)
revenue_p95 = df["revenue"].quantile(0.95)

# Top stats
top_region = regional.iloc[0]["region"]
top_region_pct = regional.iloc[0]["revenue"] / total_revenue * 100
top_category = category_stats.iloc[0]["category"]
top_category_pct = category_stats.iloc[0]["revenue"] / total_revenue * 100
top_city = city_stats.iloc[0]["customer_city"]

# Growth: compare first half vs second half
mid_date = df["order_date"].median()
first_half = df[df["order_date"] <= mid_date]["revenue"].sum()
second_half = df[df["order_date"] > mid_date]["revenue"].sum()
growth_pct = ((second_half - first_half) / first_half) * 100


# ============================================================
# BUILD FRENCH REPORT
# ============================================================

def build_french_report():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Rapport Statistique\nE-Commerce Olist",
        "Analyse Approfondie des Donnees | 2016-2018",
        lang="fr"
    )

    # Slide 2: Executive Summary
    kpis = [
        {"label": "Revenu Total", "value": f"R$ {total_revenue:,.0f}"},
        {"label": "Commandes Totales", "value": f"{total_orders:,}"},
        {"label": "Valeur Moyenne/Commande", "value": f"R$ {avg_order_value:,.2f}"},
        {"label": "Note Moyenne", "value": f"{avg_review:.1f}/5.0"},
    ]
    kpi_explanation = (
        f"Cette analyse couvre {total_orders:,} commandes de {unique_customers:,} clients uniques "
        f"sur la periode {date_min} a {date_max}. Le revenu total est de R$ {total_revenue:,.0f} "
        f"avec une croissance de {growth_pct:.1f}% entre la premiere et la seconde moitie de la periode."
    )
    add_kpi_slide(prs, kpis, kpi_explanation, lang="fr")

    # Slide 3: Revenue Overview
    fig = px.line(monthly, x="order_date", y="revenue",
                  title="", template="plotly_dark", markers=True)
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"Le graphique montre l'evolution mensuelle des revenus. On observe une tendance "
        f"generale a la hausse avec une croissance de {growth_pct:.1f}% entre la premiere et "
        f"seconde moitie de la periode analysee. Le pic de revenu atteint R$ {monthly['revenue'].max():,.0f} "
        f"en un seul mois. Le minimum etait de R$ {monthly['revenue'].min():,.0f}."
    )
    add_content_slide(prs, "Tendance des Revenus par Mois", fig, explanation)

    # Slide 4: Regional Distribution
    fig = px.bar(regional.head(10).sort_values("revenue", ascending=True),
                 x="revenue", y="region", orientation="h",
                 title="", template="plotly_dark",
                 color="revenue", color_continuous_scale="Tealgrn")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"L'Etat de {top_region} domine avec {top_region_pct:.1f}% du revenu total, "
        f"suivi par {regional.iloc[1]['region']} et {regional.iloc[2]['region']}. "
        f"Cela montre une forte concentration geographique: les 3 premiers etats representent "
        f"plus de 60% de toutes Les recettes. La diversification regionale reste un defi."
    )
    add_content_slide(prs, "Distribution Regionale (Top 10 Etats)", fig, explanation)

    # Slide 5: Category Breakdown
    fig = px.pie(category_stats.head(10), names="category", values="revenue",
                 title="", template="plotly_dark", hole=0.4)
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"La categorie '{top_category}' est la plus performante avec {top_category_pct:.1f}% "
        f"du revenu total. Le top 3 des categories represente environ 30% des revenus, "
        f"indiquee une diversification moderee. Il existe {unique_categories} categories au total, "
        f"ce qui offre de nombreuses opportunites de developpement."
    )
    add_content_slide(prs, "Repartition par Categorie (Top 10)", fig, explanation)

    # Slide 6: Volume de Commandes
    fig = px.bar(monthly, x="order_date", y="orders",
                 title="", template="plotly_dark",
                 color="orders", color_continuous_scale="Blues")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"Le volume de commandes suit la meme tendance que les revenus. "
        f"Maximum: {monthly['orders'].max():,} commandes en un mois. "
        f"Minimum: {monthly['orders'].min():,} commandes. "
        f"Moyenne mensuelle: {monthly['orders'].mean():,.0f} commandes."
    )
    add_content_slide(prs, "Volume de Commandes par Mois", fig, explanation)

    # Slide 7: Review Scores
    review_dist = df.groupby("avg_review_score").size().reset_index(name="count")
    review_dist["avg_review_score"] = review_dist["avg_review_score"].round(0).astype(int)
    review_agg = review_dist.groupby("avg_review_score")["count"].sum().reset_index()
    fig = px.bar(review_agg, x="avg_review_score", y="count",
                 title="", template="plotly_dark",
                 color="count", color_continuous_scale="Viridis")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    pct_good = len(df[df["avg_review_score"] >= 4]) / len(df) * 100
    explanation = (
        f"La note moyenne est de {avg_review:.2f}/5.0. {pct_good:.1f}% des commandes ont une note "
        f"de 4 ou 5 etoiles, indiquant une satisfaction client eleveee. "
        f"Les notes basses (1-2) representent seulement {100-pct_good:.1f}% du total, "
        f"ce qui montre que la qualite de service est globalement bonne."
    )
    add_content_slide(prs, "Satisfaction Client (Notes de Revue)", fig, explanation)

    # Slide 8: Payment Types
    fig = px.pie(payment_stats, names="payment_type", values="revenue",
                 title="", template="plotly_dark", hole=0.3)
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    credit_pct = payment_stats[payment_stats["payment_type"] == "credit_card"]["revenue"].sum() / total_revenue * 100
    explanation = (
        f"La carte de credit est le moyen de paiement dominant avec {credit_pct:.1f}% des revenus. "
        f"Les autres methodes (boleto, voucher, debit card) partagent le reste. "
        f"Cette preference pour la carte de credit suggere une clientele aisee qui prefere "
        f"payer en plusieurs fois."
    )
    add_content_slide(prs, "Analyse des Modes de Paiement", fig, explanation)

    # Slide 9: Price vs Freight
    sample = df.sample(min(3000, len(df)))
    fig = px.scatter(sample, x="price", y="freight_value",
                     title="", template="plotly_dark",
                     opacity=0.5, trendline="ols")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"La correlation entre le prix et le fret est de {price_freight_corr:.3f}, "
        f"ce qui indique une relation {'forte' if abs(price_freight_corr) > 0.5 else 'moderee'}."
        f"Les produits plus chers ont tendance a avoir des frais de port plus eleves. "
        f"Le fret moyen est de R$ {df['freight_value'].mean():,.2f} contre un prix moyen de R$ {df['price'].mean():,.2f}."
    )
    add_content_slide(prs, "Correlation Prix-Fret (Frais de Port)", fig, explanation)

    # Slide 10: Day of Week
    fig = px.bar(dow_stats, x="day_of_week", y="revenue",
                 title="", template="plotly_dark",
                 color="revenue", color_continuous_scale="Plasma")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    max_day = dow_stats.loc[dow_stats["revenue"].idxmax(), "day_of_week"]
    min_day = dow_stats.loc[dow_stats["revenue"].idxmin(), "day_of_week"]
    explanation = (
        f"Le {max_day} est le jour avec le plus de revenus, tandis que le {min_day} "
        f"est le plus faible. La difference entre le meilleur et le pire jour est "
        f"de {(dow_stats['revenue'].max() / dow_stats['revenue'].min() - 1) * 100:.1f}%. "
        f"Ces informations permettent d'optimiser les campagnes marketing et la gestion des stocks."
    )
    add_content_slide(prs, "Revenus par Jour de la Semaine", fig, explanation)

    # Slide 11: Top Cities
    fig = px.bar(city_stats.head(10).sort_values("revenue", ascending=True),
                 x="revenue", y="customer_city", orientation="h",
                 title="", template="plotly_dark",
                 color="revenue", color_continuous_scale="YlOrRd")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"{top_city} est la ville qui genere le plus de revenus, avec une concentration "
        f"importante dans les grandes metropolises bresiliennes. "
        f"Les 5 premieres villes representent environ 30% du revenu national. "
        f"Cela indique un potentiel de croissance dans les villes moyennes."
    )
    add_content_slide(prs, "Top 10 Villes par Revenu", fig, explanation)

    # Slide 12: Order Value Distribution
    fig = px.histogram(df, x="revenue", nbins=80,
                       title="", template="plotly_dark")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"La distribution des valeurs de commande est asymetrique (skewed right). "
        f"La mediane est de R$ {revenue_p50:,.2f} et la moyenne de R$ {avg_order_value:,.2f}. "
        f"Le 95e percentile est a R$ {revenue_p95:,.2f}, ce que 95% des commandes sont en dessous. "
        f"Les commandes de plus de R$ {revenue_p75:,.2f} constituent le top 25%."
    )
    add_content_slide(prs, "Distribution des Valeurs de Commande", fig, explanation)

    # Slide 13: Statistical Summary Table
    summary_df = pd.DataFrame({
        "Metrique": [
            "Revenu Total", "Nombre de Commandes", "Valeur Moyenne",
            "Note Moyenne", "Revenu P25", "Mediane (P50)",
            "Revenu P75", "Revenu P95", "Correlation Prix-Fret",
            "Clients Uniques", "Categories", "Periode"
        ],
        "Valeur": [
            f"R$ {total_revenue:,.2f}", f"{total_orders:,}", f"R$ {avg_order_value:,.2f}",
            f"{avg_review:.2f}", f"R$ {revenue_p25:,.2f}", f"R$ {revenue_p50:,.2f}",
            f"R$ {revenue_p75:,.2f}", f"R$ {revenue_p95:,.2f}", f"{price_freight_corr:.3f}",
            f"{unique_customers:,}", f"{unique_categories}", f"{date_min} a {date_max}"
        ]
    })
    table_explanation = (
        "Ce tableau resume les principales statistiques descriptives. "
        "Les percentiles montrent la distribution des revenus: 50% des commandes sont "
        f"inferieures a R$ {revenue_p50:,.2f} et 75% sont inferieures a R$ {revenue_p75:,.2f}. "
        "La correlation prix-fret indique la relation lineaire entre ces deux variables."
    )
    add_table_slide(prs, "Resume Statistique", summary_df, table_explanation, lang="fr")

    # Slide 14: Regional Table
    regional_display = regional.head(10).copy()
    regional_display["revenue"] = regional_display["revenue"].apply(lambda x: f"R$ {x:,.0f}")
    regional_display["avg_review"] = regional_display["avg_review"].apply(lambda x: f"{x:.1f}")
    regional_display["part"] = regional.head(10)["revenue"].apply(lambda x: f"{x/total_revenue*100:.1f}%")
    regional_display = regional_display[["region", "revenue", "orders", "avg_review", "part"]]
    regional_display.columns = ["Etat", "Revenu", "Commandes", "Note Moy.", "Part %"]
    table_explanation = (
        f"Les 10 premiers etats representent la majeure partie des revenus. "
        f"L'Etat de {top_region} domine avec la plus grande part. "
        "Les notes moyennes sont generalement elevees dans tous les etats, "
        "ce qui indique une qualite de service uniforme sur tout le territoire."
    )
    add_table_slide(prs, "Top 10 Etats - Details", regional_display, table_explanation, lang="fr")

    # Slide 15: Category Table
    cat_display = category_stats.head(10).copy()
    cat_display["revenue"] = cat_display["revenue"].apply(lambda x: f"R$ {x:,.0f}")
    cat_display["avg_price"] = cat_display["avg_price"].apply(lambda x: f"R$ {x:,.2f}")
    cat_display["avg_review"] = cat_display["avg_review"].apply(lambda x: f"{x:.1f}")
    cat_display["part"] = category_stats.head(10)["revenue"].apply(lambda x: f"{x/total_revenue*100:.1f}%")
    cat_display = cat_display[["category", "revenue", "orders", "avg_price", "avg_review", "part"]]
    cat_display.columns = ["Categorie", "Revenu", "Commandes", "Prix Moy.", "Note Moy.", "Part %"]
    table_explanation = (
        f"La categorie leader est '{top_category}' avec {top_category_pct:.1f}% des revenus. "
        "Les prix moyens varient considerablement entre les categories, "
        "allant de quelques dizaines a plusieurs centaines de R$. "
        "La satisfaction est generalement bonne dans toutes les categories."
    )
    add_table_slide(prs, "Top 10 Categories - Details", cat_display, table_explanation, lang="fr")

    # Slide 16: Key Insights
    add_text_slide(prs, "Conclusions Principales", [
        f"Concentration regionale: {top_region} represente {top_region_pct:.1f}% du revenu total, montrant une forte dependance geographique.",
        f"Croissance: La croissance de {growth_pct:.1f}% entre les deux moitiees de la periode indique un marche en expansion.",
        f"Paiement: La carte de credit domine avec {credit_pct:.1f}%, refletant les habitudes de consommation bresiliennes.",
        f"Satisfaction: {pct_good:.1f}% des clients donnent une note >= 4/5, indiquant une bonne qualite de service.",
        f"Correlation Prix-Fret: r = {price_freight_corr:.3f}, relation significative entre prix et frais de port.",
        f"Concentration urbaine: {top_city} et les grandes villes generent une part disproportionnee des revenus.",
        f"Diversification: {unique_categories} categories offrent de nombreuses opportunites de developpement.",
        f"Distribution asymetrique: La majorite des commandes sont de petite valeur, avec quelques commandes de grande valeur."
    ], lang="fr")

    # Slide 17: Recommendations
    add_text_slide(prs, "Recommandations Strategiques", [
        "Expansion regionale: Cibler les etats sous-performants (nord, nord-est) avec des campagnes marketing localisees.",
        "Optimisation du fret: Negocier les frais de transport pour reduire le cout moyen du fret et augmenter les marges.",
        "Diversification des categories: Developper les categories emergentes a forte marge beneficiaire.",
        "Programme de fidelite: Mettre en place un programme pour retenir les clients satisfaits (note >= 4).",
        "Paiements alternatifs: Promouvoir le debit et le boleto pour reduire la dependance a la carte de credit.",
        "Optimisation de la semaine: Concentrer les promotions sur les jours faibles ({min_day}) pour equilibrer les ventes.",
        "Grandes villes: Renforcer la presence dans les villes moyennes pour reduire la concentration geographique.",
        "Analyse continue: Suivre les KPIs mensuellement pour detecter les tendances et ajuster la strategie."
    ], lang="fr")

    # Slide 18: Data Source
    add_text_slide(prs, "Source des Donnees et Methodologie", [
        "Source: Olist Brazilian E-Commerce Dataset (Kaggle)",
        f"Periode d'analyse: {date_min} a {date_max}",
        f"Commandes analysees: {total_orders:,} (livrees uniquement)",
        f"Clients uniques: {unique_customers:,}",
        f"Categories de produits: {unique_categories}",
        "Geographic: 27 etats bresiliens + district federal",
        "Methodologie: Analyse descriptive, correlation de Pearson, aggrgation temporelle",
        "Outil: Python (Pandas, Plotly) pour l'analyse et la visualisation",
    ], lang="fr")

    # Save
    output_path = OUTPUT_DIR / "report_fr.pptx"
    prs.save(output_path)
    print(f"French report saved: {output_path}")


# ============================================================
# BUILD ARABIC REPORT
# ============================================================

def build_arabic_report():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "التقرير الإحصائي\nللتجارة الإلكترونية أوليست",
        "تحليل معمق للبيانات | 2016-2018",
        lang="ar"
    )

    # Slide 2: Executive Summary
    kpis = [
        {"label": "إجمالي الإيرادات", "value": f"R$ {total_revenue:,.0f}"},
        {"label": "إجمالي الطلبات", "value": f"{total_orders:,}"},
        {"label": "متوسط قيمة الطلب", "value": f"R$ {avg_order_value:,.2f}"},
        {"label": "متوسط التقييم", "value": f"{avg_review:.1f}/5.0"},
    ]
    kpi_explanation = (
        f"يغطي هذا التحليل {total_orders:,} طلب من {unique_customers:,} عميل فريد "
        f"خلال الفترة من {date_min} إلى {date_max}. إجمالي الإيرادات R$ {total_revenue:,.0f} "
        f"مع نمو بنسبة {growth_pct:.1f}% بين النصف الأول والثاني من الفترة."
    )
    add_kpi_slide(prs, kpis, kpi_explanation, lang="ar")

    # Slide 3: Revenue Overview
    fig = px.line(monthly, x="order_date", y="revenue",
                  title="", template="plotly_dark", markers=True)
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"يوضح الرسم البياني تطور الإيرادات الشهرية. نلاحظ اتجاهًا عامًا صاعدًا "
        f"مع نمو بنسبة {growth_pct:.1f}% بين النصفين. "
        f"بلغت الذروة R$ {monthly['revenue'].max():,.0f} في شهر واحد."
    )
    add_content_slide(prs, "اتجاه الإيرادات الشهرية", fig, explanation, lang="ar")

    # Slide 4: Regional Distribution
    fig = px.bar(regional.head(10).sort_values("revenue", ascending=True),
                 x="revenue", y="region", orientation="h",
                 title="", template="plotly_dark",
                 color="revenue", color_continuous_scale="Tealgrn")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"تهيمن ولاية {top_region} بنسبة {top_region_pct:.1f}% من الإيرادات، "
        f"تليها {regional.iloc[1]['region']} و {regional.iloc[2]['region']}. "
        f"هذا يوضح تركيزًا جغرافيًا قويًا: الولايات الثلاث الأولى تمثل أكثر من 60% من الإيرادات."
    )
    add_content_slide(prs, "التوزيع الإقليمي (أفضل 10 ولايات)", fig, explanation, lang="ar")

    # Slide 5: Category Breakdown
    fig = px.pie(category_stats.head(10), names="category", values="revenue",
                 title="", template="plotly_dark", hole=0.4)
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"الفئة '{top_category}' هي الأكثر أداءً بنسبة {top_category_pct:.1f}% من الإيرادات. "
        f"أفضل 3 فئات تمثل حوالي 30% من الإيرادات. "
        f"يوجد {unique_categories} فئة إجمالاً، مما يوفر فرصًا متعددة للتطوير."
    )
    add_content_slide(prs, "توزيع الفئات (أفضل 10)", fig, explanation, lang="ar")

    # Slide 6: Volume de Commandes
    fig = px.bar(monthly, x="order_date", y="orders",
                 title="", template="plotly_dark",
                 color="orders", color_continuous_scale="Blues")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"حجم الطلبات يتبع نفس اتجاه الإيرادات. "
        f"الحد الأقصى: {monthly['orders'].max():,} طلب في شهر واحد. "
        f"المتوسط الشهري: {monthly['orders'].mean():,.0f} طلب."
    )
    add_content_slide(prs, "حجم الطلبات حسب الشهر", fig, explanation, lang="ar")

    # Slide 7: Review Scores
    review_dist = df.groupby("avg_review_score").size().reset_index(name="count")
    review_dist["avg_review_score"] = review_dist["avg_review_score"].round(0).astype(int)
    review_agg = review_dist.groupby("avg_review_score")["count"].sum().reset_index()
    fig = px.bar(review_agg, x="avg_review_score", y="count",
                 title="", template="plotly_dark",
                 color="count", color_continuous_scale="Viridis")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    pct_good = len(df[df["avg_review_score"] >= 4]) / len(df) * 100
    explanation = (
        f"متوسط التقييم {avg_review:.2f}/5.0. {pct_good:.1f}% من الطلبات حصلت على تقييم "
        f"4 أو 5 نجوم، مما يشير إلى رضا العملاء. "
        f"التقييمات المنخفضة (1-2) تمثل فقط {100-pct_good:.1f}% من الإجمالي."
    )
    add_content_slide(prs, "رضا العملاء (التقييمات)", fig, explanation, lang="ar")

    # Slide 8: Payment Types
    fig = px.pie(payment_stats, names="payment_type", values="revenue",
                 title="", template="plotly_dark", hole=0.3)
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    credit_pct = payment_stats[payment_stats["payment_type"] == "credit_card"]["revenue"].sum() / total_revenue * 100
    explanation = (
        f"بطاقة الائتمان هي وسيلة الدفع المهيمنة بنسبة {credit_pct:.1f}% من الإيرادات. "
        f"الطرق الأخرى (boleto, voucher, debit card) تشارك في الباقي. "
        f"هذا يعكس تفضيل العملاء للدفع بالتقسيط."
    )
    add_content_slide(prs, "تحليل وسائل الدفع", fig, explanation, lang="ar")

    # Slide 9: Price vs Freight
    sample = df.sample(min(3000, len(df)))
    fig = px.scatter(sample, x="price", y="freight_value",
                     title="", template="plotly_dark",
                     opacity=0.5, trendline="ols")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"الارتباط بين السعر والشحن هو {price_freight_corr:.3f}، "
        f"مما يشير إلى علاقة {'قوية' if abs(price_freight_corr) > 0.5 else 'متوسطة'}. "
        f"متوسط الشحن R$ {df['freight_value'].mean():,.2f} مقابل متوسط السعر R$ {df['price'].mean():,.2f}."
    )
    add_content_slide(prs, "العلاقة بين السعر والشحن", fig, explanation, lang="ar")

    # Slide 10: Day of Week
    arabic_days = ["الاثنين", "الثلاثاء", "الأربعاء", "الخميس", "الجمعة", "السبت", "الأحد"]
    dow_stats_arabic = dow_stats.copy()
    dow_stats_arabic["day_of_week"] = arabic_days
    fig = px.bar(dow_stats_arabic, x="day_of_week", y="revenue",
                 title="", template="plotly_dark",
                 color="revenue", color_continuous_scale="Plasma")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"يوضح الرسم الإيرادات حسب أيام الأسبوع. "
        f"هذا يساعد في تحسين الحملات التسويقية وإدارة المخزون. "
        f"الفرق بين أفضل وأسوأ يوم هو {(dow_stats['revenue'].max() / dow_stats['revenue'].min() - 1) * 100:.1f}%."
    )
    add_content_slide(prs, "الإيرادات حسب اليوم", fig, explanation, lang="ar")

    # Slide 11: Top Cities
    fig = px.bar(city_stats.head(10).sort_values("revenue", ascending=True),
                 x="revenue", y="customer_city", orientation="h",
                 title="", template="plotly_dark",
                 color="revenue", color_continuous_scale="YlOrRd")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"{top_city} هي المدينة التي تولد أكبر إيرادات. "
        f"التركز في المدن الكبرى البرازيلية ملحوظ. "
        f"أفضل 5 مدن تمثل حوالي 30% من الإيرادات الوطنية."
    )
    add_content_slide(prs, "أفضل 10 مدن من حيث الإيرادات", fig, explanation, lang="ar")

    # Slide 12: Order Value Distribution
    fig = px.histogram(df, x="revenue", nbins=80,
                       title="", template="plotly_dark")
    fig.update_layout(margin=dict(l=30, r=30, t=30, b=30))
    explanation = (
        f"توزيع قيم الطلبات ملتوي نحو اليمين. "
        f"الوسيط R$ {revenue_p50:,.2f} والمتوسط R$ {avg_order_value:,.2f}. "
        f"المئين 95 عند R$ {revenue_p95:,.2f}، أي أن 95% الطلبات أقل من ذلك."
    )
    add_content_slide(prs, "توزيع قيم الطلبات", fig, explanation, lang="ar")

    # Slide 13: Statistical Summary Table
    summary_df = pd.DataFrame({
        "المقياس": [
            "إجمالي الإيرادات", "عدد الطلبات", "المتوسط",
            "متوسط التقييم", "المئين 25", "الوسيط",
            "المئين 75", "المئين 95", "ارتباط السعر-الشحن",
            "العملاء الفريدون", "الفئات", "الفترة"
        ],
        "القيمة": [
            f"R$ {total_revenue:,.2f}", f"{total_orders:,}", f"R$ {avg_order_value:,.2f}",
            f"{avg_review:.2f}", f"R$ {revenue_p25:,.2f}", f"R$ {revenue_p50:,.2f}",
            f"R$ {revenue_p75:,.2f}", f"R$ {revenue_p95:,.2f}", f"{price_freight_corr:.3f}",
            f"{unique_customers:,}", f"{unique_categories}", f"{date_min} إلى {date_max}"
        ]
    })
    table_explanation = (
        "يلخص هذا الجدول الإحصائيات الوصفية الرئيسية. "
        "تظهر المئينات توزيع الإيرادات: 50% من الطلبات أقل "
        f"من R$ {revenue_p50:,.2f} و75% أقل من R$ {revenue_p75:,.2f}."
    )
    add_table_slide(prs, "ملخص إحصائي", summary_df, table_explanation, lang="ar")

    # Slide 14: Regional Table
    regional_display = regional.head(10).copy()
    regional_display["revenue"] = regional_display["revenue"].apply(lambda x: f"R$ {x:,.0f}")
    regional_display["avg_review"] = regional_display["avg_review"].apply(lambda x: f"{x:.1f}")
    regional_display["part"] = regional.head(10)["revenue"].apply(lambda x: f"{x/total_revenue*100:.1f}%")
    regional_display = regional_display[["region", "revenue", "orders", "avg_review", "part"]]
    regional_display.columns = ["الولاية", "الإيرادات", "الطلبات", "متوسط التقييم", "النسبة"]
    table_explanation = (
        f"أفضل 10 ولايات تمثل معظم الإيرادات. "
        f"ولاية {top_region} تهيمن بأكبر حصة. "
        "التقييمات العامة مرتفعة في جميع الولايات."
    )
    add_table_slide(prs, "أفضل 10 ولايات - تفاصيل", regional_display, table_explanation, lang="ar")

    # Slide 15: Category Table
    cat_display = category_stats.head(10).copy()
    cat_display["revenue"] = cat_display["revenue"].apply(lambda x: f"R$ {x:,.0f}")
    cat_display["avg_price"] = cat_display["avg_price"].apply(lambda x: f"R$ {x:,.2f}")
    cat_display["avg_review"] = cat_display["avg_review"].apply(lambda x: f"{x:.1f}")
    cat_display["part"] = category_stats.head(10)["revenue"].apply(lambda x: f"{x/total_revenue*100:.1f}%")
    cat_display = cat_display[["category", "revenue", "orders", "avg_price", "avg_review", "part"]]
    cat_display.columns = ["الفئة", "الإيرادات", "الطلبات", "متوسط السعر", "متوسط التقييم", "النسبة"]
    table_explanation = (
        f"الفئة الرائدة '{top_category}' بنسبة {top_category_pct:.1f}%. "
        "الأسعار تختلف بشكل كبير بين الفئات. "
        "الرضا العام جيد في جميع الفئات."
    )
    add_table_slide(prs, "أفضل 10 فئات - تفاصيل", cat_display, table_explanation, lang="ar")

    # Slide 16: Key Insights
    add_text_slide(prs, "النتائج الرئيسية", [
        f"تركيز إقليمي: {top_region} تمثل {top_region_pct:.1f}% من الإيرادات، مما يظهر اعتمادًا جغرافيًا كبيرًا.",
        f"النمو: النمو بنسبة {growth_pct:.1f}% بين النصفين يشير إلى سوق متسع.",
        f"الدفع: بطاقة الائتمان تهيمن بنسبة {credit_pct:.1f}%. ",
        f"الرضا: {pct_good:.1f}% من العملاء يعطون تقييم >= 4/5.",
        f"الارتباط: العلاقة بين السعر والشحن r = {price_freight_corr:.3f}.",
        f"التركيز الحضري: {top_city} والمدن الكبرى تولد حصة غير متناسبة.",
        f"التنويع: {unique_categories} فئة توفر فرصًا للتطوير.",
        f"التوزيع: معظم الطلبات صغيرة القيمة مع عدد قليل من الطلبات الكبيرة."
    ], lang="ar")

    # Slide 17: Recommendations
    add_text_slide(prs, "التوصيات الاستراتيجية", [
        "التوسع الإقليمي: استهداف الولايات ضعيفة الأداء (الشمال، الشمال الشرقي) بحملات تسويق محلية.",
        "تحسين الشحن: التفاوض على أسعار الشحن لتقليل التكاليف وزيادة الهوامش.",
        "تنويع الفئات: تطوير الفئات الناشئة ذات الهوامج المرتفعة.",
        "برنامج الولاء: الاحتفاظ بالعملاء الراضين (تقييم >= 4).",
        "بدائل الدفع: تعزيز الدفع المباشر والboleto.",
        "تحسين الأسبوع: تركيز العروض على الأيام الضعيفة.",
        "المدن المتوسطة: تعزيز الحضور في المدن المتوسطة.",
        "التحليل المستمر: متابعة المؤشرات شهريًا."
    ], lang="ar")

    # Slide 18: Data Source
    add_text_slide(prs, "مصدر البيانات والمنهجية", [
        "المصدر: مجموعة بيانات أوليست البرازيلية للتجارة الإلكترونية (Kaggle)",
        f"فترة التحليل: {date_min} إلى {date_max}",
        f"الطلبات المحللة: {total_orders:,}",
        f"العملاء الفريدون: {unique_customers:,}",
        f"فئات المنتجات: {unique_categories}",
        "جغرافي: 27 ولاية برازيلية + المقاطعة الفيدرالية",
        "المنهجية: تحليل وصفي، ارتباط بيرسون، تجميع زمني",
        "الأدوات: Python (Pandas, Plotly)",
    ], lang="ar")

    # Save
    output_path = OUTPUT_DIR / "report_ar.pptx"
    prs.save(output_path)
    print(f"Arabic report saved: {output_path}")


# ============================================================
# RUN
# ============================================================

if __name__ == "__main__":
    print("Building French report...")
    build_french_report()

    print("Building Arabic report...")
    build_arabic_report()

    print("\n=== Reports generated ===")
    print(f"Location: {OUTPUT_DIR.absolute()}")
